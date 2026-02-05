import os
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class SlideService:
    """Parse PDF and PPTX files into a structured outline with extracted images."""

    # Maximum images to keep per lecture after filtering.
    MAX_IMAGES = 10
    # Minimum pixel area for an image to be considered (filters out icons).
    MIN_PIXEL_AREA = 100 * 100  # 100 x 100 px

    @staticmethod
    def parse(file_path: str, session_data_dir: str) -> dict:
        """Dispatch to the correct parser based on file extension.

        Returns::

            {
                "sections": [
                    {"title": str, "bullets": [str], "slide_refs": [int], "image_paths": [str]},
                    ...
                ],
                "raw_pages": [{"page": int, "text": str}, ...],
            }
        """
        images_dir = os.path.join(session_data_dir, "slides", "images")
        os.makedirs(images_dir, exist_ok=True)

        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return SlideService._parse_pdf(file_path, images_dir)
        elif ext in (".pptx", ".ppt"):
            return SlideService._parse_pptx(file_path, images_dir)
        else:
            raise ValueError(f"Unsupported slide format: {ext}")

    # ------------------------------------------------------------------
    # PDF parsing
    # ------------------------------------------------------------------
    @staticmethod
    def _parse_pdf(file_path: str, images_dir: str) -> dict:
        doc = fitz.open(file_path)
        sections: list[dict] = []
        raw_pages: list[dict] = []
        all_candidate_images: list[dict] = []

        for page_idx in range(len(doc)):
            page = doc[page_idx]
            text = page.get_text().strip()
            raw_pages.append({"page": page_idx + 1, "text": text})

            # Derive section title from the first non-empty line
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            title = lines[0] if lines else f"Page {page_idx + 1}"
            bullets = lines[1:] if len(lines) > 1 else []

            # Extract embedded images
            image_list = page.get_images(full=True)
            for img_tuple in image_list:
                xref = img_tuple[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    width, height = base_image["width"], base_image["height"]
                    if width * height < SlideService.MIN_PIXEL_AREA:
                        continue
                    ext = base_image.get("ext", "png")
                    filename = f"slide_{page_idx + 1:02d}_img_{xref}.{ext}"
                    filepath = os.path.join(images_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img_bytes)
                    all_candidate_images.append({
                        "page": page_idx + 1,
                        "path": filepath,
                        "filename": filename,
                        "area": width * height,
                    })
                except Exception:
                    continue  # skip corrupt or unreadable embedded images

            sections.append({
                "title": title,
                "bullets": bullets,
                "slide_refs": [page_idx + 1],
                "image_paths": [],  # filled after global filtering
            })

        doc.close()

        # Filter and cap images globally, then attach to sections
        selected = SlideService._filter_images(all_candidate_images)
        page_to_images: dict[int, list[str]] = {}
        for img in selected:
            page_to_images.setdefault(img["page"], []).append(img["filename"])
        for section in sections:
            page_num = section["slide_refs"][0]
            section["image_paths"] = page_to_images.get(page_num, [])

        return {"sections": sections, "raw_pages": raw_pages}

    # ------------------------------------------------------------------
    # PPTX parsing
    # ------------------------------------------------------------------
    @staticmethod
    def _parse_pptx(file_path: str, images_dir: str) -> dict:
        prs = Presentation(file_path)
        sections: list[dict] = []
        raw_pages: list[dict] = []
        all_candidate_images: list[dict] = []

        for slide_idx, slide in enumerate(prs.slides):
            title: str | None = None
            texts: list[str] = []

            for shape in slide.shapes:
                # --- Text extraction ---
                if shape.has_text_frame:
                    frame_text = shape.text_frame.text.strip()
                    if not frame_text:
                        continue
                    # Title placeholder has idx == 0
                    if (
                        title is None
                        and hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    ):
                        title = frame_text
                        continue
                    texts.append(frame_text)

                # --- Image extraction (picture shapes only) ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        # Use shape dimensions (EMU) converted to approx pixels at 96 DPI
                        w_px = int(shape.width / 914400 * 96)
                        h_px = int(shape.height / 914400 * 96)
                        if w_px * h_px < SlideService.MIN_PIXEL_AREA:
                            continue
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        filename = f"slide_{slide_idx + 1:02d}_img_{shape.shape_id}.{ext}"
                        filepath = os.path.join(images_dir, filename)
                        with open(filepath, "wb") as f:
                            f.write(image.blob)
                        all_candidate_images.append({
                            "page": slide_idx + 1,
                            "path": filepath,
                            "filename": filename,
                            "area": w_px * h_px,
                        })
                    except Exception:
                        continue

            # Fallback title: first text block, or generic label
            if title is None:
                title = texts[0] if texts else f"Slide {slide_idx + 1}"
                bullets = texts[1:] if len(texts) > 1 else []
            else:
                bullets = texts

            raw_pages.append({
                "page": slide_idx + 1,
                "text": f"{title}\n" + "\n".join(bullets),
            })
            sections.append({
                "title": title,
                "bullets": bullets,
                "slide_refs": [slide_idx + 1],
                "image_paths": [],
            })

        # Filter and attach images
        selected = SlideService._filter_images(all_candidate_images)
        page_to_images: dict[int, list[str]] = {}
        for img in selected:
            page_to_images.setdefault(img["page"], []).append(img["filename"])
        for section in sections:
            page_num = section["slide_refs"][0]
            section["image_paths"] = page_to_images.get(page_num, [])

        return {"sections": sections, "raw_pages": raw_pages}

    # ------------------------------------------------------------------
    # Image filtering
    # ------------------------------------------------------------------
    @staticmethod
    def _filter_images(candidates: list[dict]) -> list[dict]:
        """Sort by pixel area descending, deduplicate, cap at MAX_IMAGES."""
        seen: set[str] = set()
        unique: list[dict] = []
        for img in sorted(candidates, key=lambda i: i["area"], reverse=True):
            if img["filename"] not in seen:
                seen.add(img["filename"])
                unique.append(img)
        return unique[: SlideService.MAX_IMAGES]
